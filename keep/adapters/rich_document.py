"""Rich Document adapter - ingests DOCX, PPTX, XLSX, EPUB, RTF, ODT."""

from __future__ import annotations

import io
import json
import logging
from datetime import datetime, timezone
from html.parser import HTMLParser
from pathlib import Path
from typing import Any, Callable
from uuid import uuid4

from sqlalchemy.ext.asyncio import AsyncSession

from app.adapters import ContentItemResult
from app.models.content_item import ContentItem
from app.models.content_segment import ContentSegment
from app.storage.base import StorageBackend

logger = logging.getLogger(__name__)


def _split_paragraphs(text: str) -> list[str]:
    parts = [part.strip() for part in text.split("\n\n") if part.strip()]
    return parts or [text.strip()]


def _extract_docx(data: bytes) -> tuple[str, int]:
    from docx import Document

    document = Document(io.BytesIO(data))
    paragraphs = [
        paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()
    ]
    text = "\n\n".join(paragraphs)
    return text, max(1, len(paragraphs) // 30)


def _extract_pptx(data: bytes) -> tuple[str, int]:
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for slide in presentation.slides:
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text.strip())
        if parts:
            slides.append("\n".join(parts))
    text = "\n\n".join(slides)
    return text, len(presentation.slides)


def _extract_xlsx(data: bytes) -> tuple[str, int]:
    from openpyxl import load_workbook

    workbook = load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    rows: list[str] = []
    for sheet in workbook.worksheets:
        for row in sheet.iter_rows(values_only=True):
            line = ", ".join(str(value) for value in row if value is not None)
            if line.strip():
                rows.append(line)
    workbook.close()
    text = "\n".join(rows)
    return text, len(rows)


class _TextExtractor(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self._parts: list[str] = []

    def handle_data(self, data: str) -> None:
        self._parts.append(data)

    def get_text(self) -> str:
        return " ".join(self._parts)


def _extract_epub(data: bytes) -> tuple[str, int]:
    import ebooklib
    from ebooklib import epub

    book = epub.read_epub(io.BytesIO(data))
    chapters: list[str] = []
    for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
        content = item.get_content()
        try:
            extractor = _TextExtractor()
            extractor.feed(content.decode("utf-8", errors="replace"))
            text = extractor.get_text().strip()
            if text:
                chapters.append(text)
        except Exception:
            continue
    return "\n\n".join(chapters), len(chapters)


def _extract_rtf(data: bytes) -> tuple[str, int]:
    from striprtf.striprtf import rtf_to_text

    text = rtf_to_text(data.decode("latin-1", errors="replace"))
    return text, 1


def _extract_odt(data: bytes) -> tuple[str, int]:
    import xml.etree.ElementTree as ET
    import zipfile

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as archive:
            with archive.open("content.xml") as content_file:
                tree = ET.parse(content_file)
        namespace = "urn:oasis:names:tc:opendocument:xmlns:text:1.0"
        paragraphs = [
            element.text or "" for element in tree.iter(f"{{{namespace}}}p") if element.text
        ]
        text = "\n\n".join(paragraph.strip() for paragraph in paragraphs if paragraph.strip())
        return text, max(1, len(paragraphs) // 30)
    except Exception as exc:
        logger.warning("ODT extraction failed: %s", exc)
        return "", 1


_EXTRACTORS: dict[str, Callable[[bytes], tuple[str, int]]] = {
    ".docx": _extract_docx,
    ".pptx": _extract_pptx,
    ".xlsx": _extract_xlsx,
    ".xls": _extract_xlsx,
    ".epub": _extract_epub,
    ".rtf": _extract_rtf,
    ".odt": _extract_odt,
}


class RichDocumentAdapter:
    source_type = "rich_document"
    display_name = "Rich Document"
    accepted_mimes = [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel",
        "application/epub+zip",
        "application/rtf",
        "text/rtf",
        "application/vnd.oasis.opendocument.text",
    ]
    accepted_extensions = [".docx", ".pptx", ".xlsx", ".xls", ".epub", ".rtf", ".odt"]
    supported_export_formats = ["txt", "md", "json"]

    async def ingest(
        self,
        *,
        db: AsyncSession,
        user_id: int,
        options: dict[str, Any],
        storage: StorageBackend,
    ) -> ContentItemResult:
        upload = options["upload"]
        title_override = options.get("title")

        data = await upload.read()
        filename = upload.filename or "document"
        ext = Path(filename).suffix.lower()
        title = (
            title_override
            or Path(filename).stem.replace("_", " ").replace("-", " ").strip()
            or filename
        )

        extractor = _EXTRACTORS.get(ext)
        if extractor is None:
            text = ""
            page_count = 1
            logger.warning("RichDocumentAdapter: no extractor for extension %s", ext)
        else:
            try:
                text, page_count = extractor(data)
            except Exception:
                logger.exception("RichDocumentAdapter: extraction failed for %s", filename)
                text = ""
                page_count = 1

        paragraphs = (
            _split_paragraphs(text)
            if text.strip()
            else ["(No text could be extracted from this file.)"]
        )
        word_count = len(text.split())
        doc_type = ext.lstrip(".")

        item_id = str(uuid4())
        now = datetime.now(timezone.utc)

        source_uri = await storage.store(
            user_id,
            item_id,
            "source/source_descriptor.json",
            json.dumps(
                {"filename": filename, "byte_size": len(data), "doc_type": doc_type}, indent=2
            ),
        )
        segments_data = [
            {"id": index, "start": 0.0, "end": 0.0, "text": paragraph, "speaker": None}
            for index, paragraph in enumerate(paragraphs)
        ]
        output_uri = await storage.store(
            user_id,
            item_id,
            "output/segments.json",
            json.dumps({"segments": segments_data, "created_at": now.isoformat()}, indent=2),
        )
        await storage.store(user_id, item_id, "output/content.txt", text)

        item = ContentItem(
            id=item_id,
            user_id=user_id,
            adapter_type=self.source_type,
            title=title,
            source_material_path=source_uri,
            source_material_mime=upload.content_type or "application/octet-stream",
            source_material_size=len(data),
            output_content_path=output_uri,
            status="ready",
            created_at=now,
            updated_at=now,
        )
        db.add(item)
        await db.flush()

        for segment in segments_data:
            db.add(
                ContentSegment(
                    content_item_id=item_id,
                    segment_index=segment["id"],
                    start_time=0.0,
                    end_time=0.0,
                    speaker=None,
                    text=segment["text"],
                )
            )
        item.properties = json.dumps(
            {
                "page_count": page_count,
                "ocr_engine": "rich_document",
                "word_count": word_count,
            }
        )
        item.item_type = "source"
        await db.flush()

        from app.services.content_service import index_content_segments

        try:
            await index_content_segments(item_id, self.source_type, db, segments=segments_data)
        except Exception:
            logger.exception("FTS indexing failed for %s", item_id)

        logger.info(
            "RichDocumentAdapter: %d paragraphs, doc_type=%s, item_id=%s",
            len(paragraphs),
            doc_type,
            item_id,
        )
        return ContentItemResult(
            content_item_id=item_id,
            title=title,
            segment_count=len(paragraphs),
        )
